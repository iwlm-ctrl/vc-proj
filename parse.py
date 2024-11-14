import json
from pptx import Presentation

def extract_ppt_text(ppt_path):
    presentation = Presentation(ppt_path)
    text_data = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_data.append(" ".join(slide_text))
    return " ".join(text_data)
